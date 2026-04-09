from __future__ import annotations

import json
import re
from collections import Counter, defaultdict
from pathlib import Path

from pypdf import PdfReader
from pptx import Presentation


ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "_generated"
OUT_DIR.mkdir(exist_ok=True)


def clean_text(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"\s+", " ", text).strip()
    return text


def extract_pdf(path: Path) -> list[str]:
    pages: list[str] = []
    reader = PdfReader(str(path))
    for page in reader.pages:
        txt = page.extract_text() or ""
        pages.append(clean_text(txt))
    return pages


def extract_pptx(path: Path) -> list[str]:
    slides: list[str] = []
    pres = Presentation(str(path))
    for slide in pres.slides:
        chunks: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                chunks.append(shape.text)
        slides.append(clean_text(" ".join(chunks)))
    return slides


def infer_topics(text_blocks: list[str]) -> list[str]:
    topic_like = []
    for block in text_blocks:
        if not block:
            continue
        first = block.split(".")[0][:160]
        if len(first.split()) >= 2:
            topic_like.append(first)
    return topic_like


def tokenize_for_freq(text: str) -> list[str]:
    words = re.findall(r"[A-Za-z][A-Za-z\-]{2,}", text.lower())
    stop = {
        "the",
        "and",
        "for",
        "that",
        "with",
        "this",
        "from",
        "are",
        "was",
        "were",
        "have",
        "has",
        "had",
        "you",
        "your",
        "their",
        "they",
        "there",
        "into",
        "using",
        "used",
        "use",
        "also",
        "can",
        "will",
        "not",
        "all",
        "any",
        "one",
        "two",
        "three",
        "when",
        "where",
        "which",
        "what",
        "why",
        "how",
        "than",
        "then",
        "each",
        "such",
        "these",
        "those",
        "only",
        "more",
        "most",
        "very",
    }
    return [w for w in words if w not in stop]


def main() -> None:
    material_files = sorted(
        [
            *ROOT.rglob("*.pdf"),
            *ROOT.rglob("*.pptx"),
        ]
    )

    records = []
    token_counter: Counter[str] = Counter()
    stem_groups: defaultdict[str, int] = defaultdict(int)

    for fp in material_files:
        rel = fp.relative_to(ROOT).as_posix()
        if fp.suffix.lower() == ".pdf":
            blocks = extract_pdf(fp)
            block_kind = "pages"
        else:
            blocks = extract_pptx(fp)
            block_kind = "slides"

        joined = "\n".join(b for b in blocks if b)
        tokens = tokenize_for_freq(joined)
        token_counter.update(tokens)

        stem = re.sub(r"^[0-9_\- ]+", "", fp.stem).strip().lower()
        stem_groups[stem] += 1

        records.append(
            {
                "file": rel,
                "type": fp.suffix.lower().lstrip("."),
                "block_kind": block_kind,
                "blocks": blocks,
                "topics_guess": infer_topics(blocks[:10]),
            }
        )

    extraction_payload = {
        "subject_guess": "Software Testing",
        "file_count": len(records),
        "records": records,
        "repeated_topic_stems": sorted(stem_groups.items(), key=lambda x: (-x[1], x[0])),
        "top_terms": token_counter.most_common(250),
    }

    (OUT_DIR / "extracted_materials.json").write_text(
        json.dumps(extraction_payload, indent=2, ensure_ascii=True), encoding="utf-8"
    )

    lines = [
        "# Extracted Materials Overview",
        "",
        f"Subject guess: {extraction_payload['subject_guess']}",
        f"Files processed: {extraction_payload['file_count']}",
        "",
        "## Files",
    ]
    for rec in records:
        non_empty = sum(1 for b in rec["blocks"] if b)
        lines.append(f"- {rec['file']} ({rec['type']}, {rec['block_kind']}: {len(rec['blocks'])}, non-empty: {non_empty})")

    lines.extend(["", "## Repeated Topic Stems"])
    for stem, cnt in extraction_payload["repeated_topic_stems"]:
        if cnt > 1:
            lines.append(f"- {stem}: {cnt}")

    lines.extend(["", "## Top Terms"])
    for term, count in extraction_payload["top_terms"][:80]:
        lines.append(f"- {term}: {count}")

    (OUT_DIR / "extraction_overview.md").write_text("\n".join(lines), encoding="utf-8")


if __name__ == "__main__":
    main()
